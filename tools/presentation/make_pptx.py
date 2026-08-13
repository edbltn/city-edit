"""Wrap the rendered panels in a .pptx — one full-bleed 16:9 slide each.

The slides ARE the images; PowerPoint is only the carrier, so nothing here
re-lays-out anything. Run after render.sh:

    uv run --with python-pptx python tools/presentation/make_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

OUT_DIR = Path(__file__).resolve().parent.parent.parent / "docs/presentation/story"
PNG_DIR = OUT_DIR / "png"
OUT = OUT_DIR / "city-edit-story.pptx"

W, H = Inches(13.333), Inches(7.5)      # 16:9 at PowerPoint's widescreen size

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
blank = prs.slide_layouts[6]            # the only layout with no placeholders

pngs = sorted(PNG_DIR.glob("*.png"))
if not pngs:
    raise SystemExit(f"no panels in {PNG_DIR} — run render.sh first")

for png in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png), 0, 0, width=W, height=H)

prs.save(OUT)
print(f"wrote {OUT} ({len(pngs)} slides, {OUT.stat().st_size / 1024:.0f} KB)")
